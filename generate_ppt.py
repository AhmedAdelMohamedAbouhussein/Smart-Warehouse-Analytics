import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()

    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Smart Warehouse System"
    subtitle.text = (
        "Optimizing Logistics through Advanced Computing Algorithms\n\n"
        "Presented by:\n"
        "- Abdelrahman mostafa (231027174)\n"
        "- Yassin eslam (231017683)\n"
        "- Ahmed adel (231007728)\n"
        "- Omar wael (231007796)"
    )

    # Helper function to add content slides
    def add_content_slide(title_text, bullet_points):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = title_text
        
        tf = slide.placeholders[1].text_frame
        tf.text = bullet_points[0]
        
        for point in bullet_points[1:]:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0

    # Slide 2: Project Overview
    add_content_slide("Project Overview & Objectives", [
        "Goal: Design an automated Warehouse Management System (WMS) for complete product lifecycle.",
        "Inbound Receiving: Efficient intake of new inventory.",
        "Slotting (Bin Packing): Optimizing product storage locations.",
        "Outbound Picking (TSP): Shortest sequence item collection.",
        "Replenishment: Automated restocking triggers.",
        "Challenge: High-velocity items must be near packing stations."
    ])

    # Slide 3: Warehouse Model
    add_content_slide("Warehouse Model", [
        "Modeled as a Spatial Graph.",
        "Nodes: Shelf locations, pick stations, and packing docks.",
        "Edges: Aisle segments with travel-time weights.",
        "Layout: 30x80 grid with clearly defined zones.",
        "[INSERT SCREENSHOT: 2D Warehouse Floor Map]"
    ])

    # Slide 4: Algorithm 1 - Dijkstra’s Algorithm
    add_content_slide("Algorithm 1 - Dijkstra’s Algorithm", [
        "Role: Shortest-path navigation.",
        "Applied to find minimum distance between any two points.",
        "Accounts for turns and aisle congestion via weights.",
        "Technical Detail: Uses a Min-Heap (Priority Queue).",
        "Complexity: O((V + E) log V) for real-time responsiveness."
    ])

    # Slide 5: Algorithm 2 - Bin Packing (Slotting)",
    add_content_slide("Algorithm 2 - Bin Packing (Slotting)", [
        "Role: Product placement strategy using Velocity.",
        "ABC Strategy:",
        " - Class A (Top 20%): High demand, closest to stations.",
        " - Class B (Next 30%): Medium demand, mid-range zones.",
        " - Class C (Remaining 50%): Low demand, far zones.",
        "FFD Approach: First-Fit Decreasing ensures shelf efficiency.",
        "[INSERT SCREENSHOT: Color-coded Slotting Plan]"
    ])

    # Slide 6: Algorithm 3 - DP Pick Route Optimization
    add_content_slide("Algorithm 3 - DP Pick Route Optimization", [
        "Role: Multi-order 'Wave Picking' optimization.",
        "The Problem: Traveling Salesperson Problem (TSP).",
        "The Solution: Held-Karp Dynamic Programming.",
        "Memoization: Uses a memo table to avoid redundant calculations.",
        "Benefit: Guaranteed mathematically optimal route."
    ])

    # Slide 7: Held-Karp vs. Greedy
    add_content_slide("Held-Karp vs. Greedy Nearest Neighbor", [
        "Greedy: Always picks next closest item (Fast but suboptimal).",
        "Held-Karp (DP): Looks at the whole trip at once.",
        "Performance Gap: DP typically reduces distance by 10%–25%.",
        "[INSERT SCREENSHOT: Pick Route Visualization]"
    ])

    # Slide 8: Algorithm 4 - Greedy Order Assignment
    add_content_slide("Algorithm 4 - Greedy Order Assignment", [
        "Role: Real-time workload balancing.",
        "Rule: Assign to picker with shortest distance to first item.",
        "Efficiency: Prevents bottlenecks and balances utilization.",
        "Real-time Logic: Triggers new picking waves if window is missed."
    ])

    # Slide 9: Replenishment & Simulation
    add_content_slide("Replenishment & Simulation", [
        "Automation: Stock < 5 units triggers replenishment.",
        "Carts follow Dijkstra-optimized paths from Receiving.",
        "Simulation: Real-time ticker with WebSocket updates.",
        "Animates multiple pickers simultaneously.",
        "[INSERT SCREENSHOT: Replenishment Log Table]"
    ])

    # Slide 10: Results & Dashboard
    add_content_slide("Results & Picker Efficiency", [
        "Total Orders Completed: Success throughput.",
        "Average Cycle Time: From creation to completion.",
        "Average Utilization %: Picking time vs. idle time.",
        "Total Walking Distance: Savings from DP optimization.",
        "[INSERT SCREENSHOT: Efficiency Dashboard Graphs]"
    ])

    # Slide 11: Conclusion
    add_content_slide("Conclusion", [
        "Efficiency Gains: Rank-based Slotting + Optimal TSP.",
        "Scalability: Handles real-time arrivals and dynamic assignments.",
        "Future Work: Machine learning for demand prediction."
    ])

    # Slide 12: Q&A
    add_content_slide("Q&A", [
        "Thank You!",
        "Questions?"
    ])

    # Save the presentation
    prs.save("Smart_Warehouse_Presentation.pptx")
    print("Presentation created successfully: Smart_Warehouse_Presentation.pptx")

if __name__ == "__main__":
    create_presentation()
